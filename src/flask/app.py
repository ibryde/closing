import json
import os
from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS

from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = './uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route('/convert', methods=['POST'])
def convert_json_to_ppt():
    if 'jsonFile' not in request.files:
        return jsonify({'success': False, 'error': 'No file part'})
    
    file = request.files['jsonFile']
    if file.filename == '':
        return jsonify({'success': False, 'error': 'No selected file'})
    
    if file:
        filename = secure_filename(file.filename)
        json_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(json_path)

        try:
            # Load JSON data
            with open(json_path, 'r') as f:
                data = json.load(f)

            # Handle repeated keys by appending "_1", "_2", etc.
            processed_data = {}
            key_counts = {}

            for key, value in data.items():
                original_key = key
                count = key_counts.get(key, 0)
                if count > 0:
                    key = f"{key}_{count}"
                key_counts[original_key] = count + 1
                processed_data[key] = value

            # Create a PowerPoint presentation
            ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], f'{os.path.splitext(filename)[0]}.pptx')
            prs = Presentation()

            # Helper function to add a slide with data
            def add_slide_with_data(prs, slide_title, data_chunk):
                slide_layout = prs.slide_layouts[1]  # Use the Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_title

                rows = len(data_chunk)
                cols = 2  # Two columns as per requirement
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)

                table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

                # Add the headers
                table.cell(0, 0).text = "Nom"
                table.cell(0, 1).text = "Quantité"

                # Add the data
                for row_num, (key, value) in enumerate(data_chunk.items(), start=1):
                    table.cell(row_num, 0).text = str(key)
                    table.cell(row_num, 1).text = str(value)

            # Split data into chunks of 10 items
            chunk_size = 10
            num_chunks = (len(processed_data) + chunk_size - 1) // chunk_size  # Calculate the number of chunks

            for i in range(num_chunks):
                start_idx = i * chunk_size
                end_idx = min(start_idx + chunk_size, len(processed_data))
                chunk_data = dict(list(processed_data.items())[start_idx:end_idx])

                # Add a slide with the chunk data
                slide_title = os.path.splitext(filename)[0]  # Use the filename without extension as slide title
                add_slide_with_data(prs, slide_title, chunk_data)

            # Save the presentation
            prs.save(ppt_path)

            return jsonify({'success': True, 'file_url': f'/download/{os.path.basename(ppt_path)}'})
        except Exception as e:
            return jsonify({'success': False, 'error': str(e)})


@app.route('/download/<filename>')
def download_file(filename):
    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=9456)
